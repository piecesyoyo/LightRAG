from pathlib import Path
import logging
import asyncio
import aiofiles
import re
import shutil
from datetime import datetime
from typing import List, Dict
import traceback
import threading
import pipmaster as pm
from fastapi import UploadFile
from lightrag import LightRAG

class DocumentManager:
    """Handles document operations and tracking"""

    def __init__(
        self,
        input_dir: str,
        supported_extensions: tuple = (
            ".txt",
            ".md",
            ".pdf",
            ".docx",
            ".pptx",
            ".xlsx",
        ),
    ):
        """
        Initialize DocumentManager
        
        Args:
            input_dir: Directory path for input documents
            supported_extensions: Tuple of supported file extensions
        """
        self.input_dir = Path(input_dir)
        self.supported_extensions = supported_extensions
        self.indexed_files = set()
        self.temp_prefix = "__tmp_"

        # Create input directory if it doesn't exist
        self.input_dir.mkdir(parents=True, exist_ok=True)

    def scan_directory_for_new_files(self) -> List[Path]:
        """
        Scan input directory for new unindexed files
        
        Returns:
            List of paths to new files
        """
        new_files = []
        for ext in self.supported_extensions:
            logging.info(f"Scanning for {ext} files in {self.input_dir}")
            for file_path in self.input_dir.rglob(f"*{ext}"):
                if file_path not in self.indexed_files:
                    new_files.append(file_path)
        return new_files

    def scan_directory(self) -> List[Path]:
        """
        Scan input directory for all supported files
        
        Returns:
            List of paths to all supported files
        """
        new_files = []
        for ext in self.supported_extensions:
            for file_path in self.input_dir.rglob(f"*{ext}"):
                new_files.append(file_path)
        return new_files

    def mark_as_indexed(self, file_path: Path):
        """
        Mark a file as indexed
        
        Args:
            file_path: Path to the indexed file
        """
        self.indexed_files.add(file_path)

    def is_supported_file(self, filename: str) -> bool:
        """
        Check if file type is supported
        
        Args:
            filename: Name of the file to check
            
        Returns:
            True if file type is supported, False otherwise
        """
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)

    async def save_temp_file(self, file: UploadFile) -> Path:
        """
        Save uploaded file to temporary location
        
        Args:
            file: FastAPI UploadFile object
            
        Returns:
            Path to saved temporary file
        """
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_filename = f"{self.temp_prefix}{timestamp}_{file.filename}"
        
        temp_path = self.input_dir / "temp" / unique_filename
        temp_path.parent.mkdir(exist_ok=True)
        
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        return temp_path


class DocumentProcessor:
    """Handles document processing and indexing operations"""

    def __init__(self, rag: LightRAG):
        """
        Initialize DocumentProcessor
        
        Args:
            rag: LightRAG instance for document processing
        """
        self.rag = rag
        self.progress_lock = threading.Lock()
        self.scan_progress: Dict = {
            "is_scanning": False,
            "current_file": "",
            "indexed_count": 0,
            "total_files": 0,
            "progress": 0,
        }

    async def process_file_content(self, file_path: Path) -> str:
        """
        Extract content from file based on its type
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted text content
        """
        content = ""
        ext = file_path.suffix.lower()

        async with aiofiles.open(file_path, "rb") as f:
            file = await f.read()

        try:
            if ext in [".txt", ".md"]:
                content = file.decode("utf-8")
            elif ext == ".pdf":
                if not pm.is_installed("pypdf2"):
                    pm.install("pypdf2")
                from PyPDF2 import PdfReader
                from io import BytesIO

                pdf_file = BytesIO(file)
                reader = PdfReader(pdf_file)
                for page in reader.pages:
                    content += page.extract_text() + "\n"
            elif ext == ".docx":
                if not pm.is_installed("docx"):
                    pm.install("docx")
                from docx import Document
                from io import BytesIO

                docx_file = BytesIO(file)
                doc = Document(docx_file)
                content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            elif ext == ".pptx":
                if not pm.is_installed("pptx"):
                    pm.install("pptx")
                from pptx import Presentation
                from io import BytesIO

                pptx_file = BytesIO(file)
                prs = Presentation(pptx_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
        except Exception as e:
            logging.error(f"Error processing file {file_path}: {str(e)}")
            logging.error(traceback.format_exc())
            raise

        return content

    async def pipeline_enqueue_file(self, file_path: Path) -> bool:
        """
        Add a file to the processing queue
        
        Args:
            file_path: Path to the file
            
        Returns:
            True if file was successfully enqueued, False otherwise
        """
        try:
            content = await self.process_file_content(file_path)
            if content:
                await self.rag.apipeline_enqueue_documents(content)
                logging.info(f"Successfully processed and enqueued file: {file_path.name}")
                return True
            logging.error(f"No content could be extracted from file: {file_path.name}")
            return False
        except Exception as e:
            logging.error(f"Error processing file {file_path.name}: {str(e)}")
            return False
        finally:
            if file_path.name.startswith("__tmp_"):
                try:
                    file_path.unlink()
                except Exception as e:
                    logging.error(f"Error deleting temporary file {file_path}: {str(e)}")

    async def pipeline_index_files(self, file_paths: List[Path]):
        """
        Index multiple files concurrently
        
        Args:
            file_paths: List of paths to files for indexing
        """
        if not file_paths:
            return
            
        try:
            if len(file_paths) == 1:
                enqueued = await self.pipeline_enqueue_file(file_paths[0])
            else:
                tasks = [self.pipeline_enqueue_file(path) for path in file_paths]
                enqueued = any(await asyncio.gather(*tasks))

            if enqueued:
                await self.rag.apipeline_process_enqueue_documents()
        except Exception as e:
            logging.error(f"Error indexing files: {str(e)}")
            logging.error(traceback.format_exc())

    async def pipeline_index_texts(self, texts: List[str]):
        """
        Index a list of texts
        
        Args:
            texts: List of texts to index
        """
        if not texts:
            return
        await self.rag.apipeline_enqueue_documents(texts)
        await self.rag.apipeline_process_enqueue_documents() 